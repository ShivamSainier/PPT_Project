from wand import image
from wand.image import Image
from wand.drawing import Drawing
from pptx import Presentation
from pptx.util import Inches,Cm,Pt

#read image

with Image(filename="nike_black.png") as img:
    img.resize(height=200,width=600,filter='undefined',blur=0)
    img.save(filename="new_logo.png")
    print("new_logo image saved")

all_images=['image1.jpg','image2.jpg','image3.jpg','image5.jpg']
s=0
for img in all_images:
    img1=Image(filename=img)
    img2=Image(filename="new_logo.png")
    img1.clone()
    img2.clone()
    with Drawing() as draw:
        draw.composite(operator='luminize',left=200,top=300,width=img2.width,height=img2.height,image=img2)
        draw(img1)
        s=s+1
        img1.save(filename=f"img{s}.jpg")
        print(f"image{s} saved")


ppt=Presentation()


def create_presentaion(layout,image,title,placeholder):
    layout=ppt.slide_layouts[layout]
    slide=ppt.slides.add_slide(layout)
    slide.shapes.add_picture(image,top=Inches(1),left=Inches(1),height=Inches(4))
    tit=slide.shapes.title
    tit.text=title
    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.name = "Times New Roman"
    title_para.font.size = Pt(30)
    title_para.font.underline = False
    slide.placeholders[1].text=placeholder

create_presentaion(0,'img1.jpg',"first slide",'first')
create_presentaion(0,'img2.jpg',"second slide",'second')
create_presentaion(0,'img3.jpg',"third slide",'third')
create_presentaion(0,'img4.jpg',"fourth slide",'forth')

ppt.save('Project.pptx')
print("ppt created")